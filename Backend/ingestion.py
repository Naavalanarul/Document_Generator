"""
Each format stores text differently:
  - PDF  → positioned drawing commands (pypdf walks pages)
  - DOCX → XML inside a zip (python-docx walks paragraph/table nodes)
  - PPTX → per-slide XML inside a zip (python-pptx walks shapes)
  - TXT/MD/CSV → plain read

After extraction, clean_text() normalises encoding noise (glyph chars,
excess whitespace) before the text hits the LLM.

For sources longer than one LLM context window, chunk_text() splits on
paragraph boundaries with overlap — keeping chunks semantically coherent
instead of cutting mid-sentence every N characters.
"""

import re
from pathlib import Path

from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

# Extensions the ingestion layer can handle — used by upload validation too.
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md", ".csv"}

# Minimum average chars/page threshold to consider a PDF as having real text
# (scanned/image-only PDFs typically yield near-zero text per page).
_MIN_CHARS_PER_PAGE = 50


# ── Shared cleanup ─────────────────────────────────────────────────────────────

def clean_text(text: str) -> str:
    """
    Normalise extraction noise common across binary formats:
    - PDF/DOCX bullet glyphs (\\x7f, \\uf0b7, etc.) → real bullet char
    - Collapse 3+ blank lines to 2
    - Strip trailing whitespace per line
    """
    text = re.sub(r"[\x7f\x80-\x9f\uf0b7\uf0a7\u25cf\u2022]", "•", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = "\n".join(line.rstrip() for line in text.splitlines())
    return text.strip()


# ── Format-specific extractors ─────────────────────────────────────────────────

def extract_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)
    pages = [page.extract_text() or "" for page in reader.pages]
    raw = "\n\n".join(p for p in pages if p.strip())

    # Detect scanned / image-only PDFs
    page_count = len(reader.pages)
    if page_count > 0 and len(raw.strip()) / page_count < _MIN_CHARS_PER_PAGE:
        raise ValueError(
            "This PDF appears to be scanned/image-based — "
            "OCR is not currently supported. "
            f"({len(raw.strip())} chars extracted from {page_count} page(s))"
        )

    return clean_text(raw)


def extract_docx(file_path: str) -> str:
    doc = DocxDocument(file_path)
    parts = []
    for para in doc.paragraphs:
        if not para.text.strip():
            continue
        if para.style.name.startswith("Heading"):
            try:
                level = int(para.style.name.split()[-1])
            except ValueError:
                level = 1
            parts.append(f"{'#' * level} {para.text}")
        else:
            parts.append(para.text)

    # Preserve table structure with markers so the LLM can recognize it
    for table in doc.tables:
        rows = []
        for row in table.rows:
            rows.append(" | ".join(c.text.strip() for c in row.cells))
        if rows:
            parts.append("[TABLE]")
            parts.extend(rows)
            parts.append("[/TABLE]")

    return clean_text("\n".join(parts))


def extract_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        indent = "  - " if para.level > 0 else ""
                        parts.append(f"{indent}{text}")
        slides.append("\n".join(parts))
    return clean_text("\n\n".join(slides))


# ── Router ─────────────────────────────────────────────────────────────────────

def extract_text(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    extractors = {
        ".pdf":  extract_pdf,
        ".docx": extract_docx,
        ".pptx": extract_pptx,
        ".txt":  lambda p: clean_text(Path(p).read_text(errors="ignore")),
        ".md":   lambda p: clean_text(Path(p).read_text(errors="ignore")),
        ".csv":  lambda p: clean_text(Path(p).read_text(errors="ignore")),
    }
    if ext not in extractors:
        raise ValueError(
            f"Unsupported file type: {ext}. "
            f"Supported: {', '.join(extractors)}"
        )
    return extractors[ext](file_path)


# ── Web extraction ─────────────────────────────────────────────────────────────

def extract_text_from_url(url: str) -> str:
    """
    Fetch readable body text from a web page using trafilatura.
    Automatically strips navbars, ads, cookie banners, and footers.
    """
    import trafilatura
    downloaded = trafilatura.fetch_url(url)
    if not downloaded:
        raise ValueError(f"Could not fetch: {url}")
    text = trafilatura.extract(downloaded, include_tables=True)
    if not text:
        raise ValueError(f"No readable content extracted from: {url}")
    return clean_text(text)


# ── Chunker ────────────────────────────────────────────────────────────────────

def chunk_text(text: str, max_chars: int = 3000, overlap: int = 300) -> list[str]:
    """
    Split long text into overlapping chunks on paragraph boundaries.

    Why paragraph boundaries?
      Splitting every N characters blindly cuts mid-sentence, confusing
      LLM summaries. Paragraph splits keep each chunk semantically coherent.

    Why overlap?
      The last `overlap` chars of each chunk are prepended to the next,
      so context at chunk edges isn't completely lost.
    """
    text = re.sub(r"\n{3,}", "\n\n", text.strip())
    if len(text) <= max_chars:
        return [text]

    paragraphs = text.split("\n\n")
    chunks, current = [], ""

    for para in paragraphs:
        if len(current) + len(para) + 2 <= max_chars:
            current += (("\n\n" if current else "") + para)
        else:
            if current:
                chunks.append(current)
            tail = current[-overlap:] if current else ""
            current = (tail + "\n\n" + para) if tail else para

    if current:
        chunks.append(current)
    return chunks
