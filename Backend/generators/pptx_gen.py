"""
Step 6: Render a PresentationPlan into a .pptx file using python-pptx.
No LLM involvement here — pure deterministic Python.
"""
from pptx import Presentation
from pptx.util import Pt
from schemas import PresentationPlan


def generate_pptx(plan: PresentationPlan, output_path: str) -> str:
    prs = Presentation()

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = plan.title
    if plan.subtitle and len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = plan.subtitle

    # Content slides
    for slide_plan in plan.slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_plan.title

        body = slide.placeholders[1].text_frame
        body.clear()

        if slide_plan.bullets:
            body.text = slide_plan.bullets[0]
            for bullet_text in slide_plan.bullets[1:]:
                p = body.add_paragraph()
                p.text = bullet_text
                p.level = 0

        if slide_plan.speaker_notes:
            slide.notes_slide.notes_text_frame.text = slide_plan.speaker_notes

    prs.save(output_path)
    return output_path
